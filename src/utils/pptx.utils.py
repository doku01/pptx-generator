import copy
import io

import pptx
from PIL import Image


# タイトルスライド挿入
def ppt_title(prs, slide_layout, title="タイトル", name="名前"):
    prs.slides.add_slide(slide_layout)
    prs.slides[-1].placeholders[0].text = title
    prs.slides[-1].placeholders[1].text = name
    return prs


# check_placeholdes
def layout_placeholders(prs, layout_num):
    temp = copy.deepcopy(prs)
    slide_layout = temp.slide_layouts[layout_num]
    temp.slides.add_slide(slide_layout)
    for shape in temp.slides[-1].placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


# get_size_list
def get_size_list(shape):
    top = shape.top
    left = shape.left
    height = shape.height
    width = shape.width
    size_list = [top, left, height, width]
    return size_list


# プレースホルダサイズ取得
def get_size_placeholders(prs, slide):
    temp = copy.deepcopy(prs)
    shape = temp.slides[slide].placeholders[1]
    top, left, height, width = shape.top, shape.left, shape.height, shape.width
    size_list = [top, left, height, width]
    return size_list


# メインスライド挿入
def ppt_main(prs, slide_layout, title="タイトル"):
    prs.slides.add_slide(slide_layout)
    prs.slides[-1].placeholders[0].text = title
    return prs


# プレースホルダにテキスト挿入
def add_txt(prs, level, size, bold=False, text="main", font_name=False, layout=1):
    tf = prs.slides[-1].placeholders[layout].text_frame
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    font = run.font
    if font_name:
        font.name = font_name
    font.size = pptx.util.Pt(size)
    font.bold = bold
    return prs


# プレースホルダに最初の段落にテキスト挿入
def add_txt_f(prs, level, size, bold=False, text="main", layout=1):
    tf = prs.slides[-1].placeholders[layout].text_frame
    tf.text = text
    tf.paragraphs[0].font.size = pptx.util.Pt(size)  # font size
    tf.paragraphs[0].font.bold = bold  # font bold
    return prs


# 複数段落分割
def split_txts(txts):
    # txt = ["level0",0,"level1",1]
    txt_list = [txts[idx:idx + 2] for idx in range(0, len(txts), 2)]
    txt_list = [i + [40, False] for i in txt_list]
    return txt_list


# 複数行入力
def add_txt_multi(prs, txt_list, font_name=False):
    for i in txt_list:
        text, level, size, bold = i
        prs = add_txt(prs, level=level, size=size, text=text, bold=bold, font_name=font_name, layout=1)
    return prs


# 画像インメモリ保存
def save_memory_PIL(PIL_image):
    item = io.BytesIO()
    PIL_image.save(item, "png")
    item.seek(0)
    return item


# 横長アスペクト比チェック
def aspect_yokonaga(width, height):
    result = True
    if int(height / width):
        result = False
    return result, height / width


# プレースホルダ縮小
def shrink_size(size_list, h_mag=1, w_mag=1):
    top, left, height, width = [int(i) for i in size_list]
    height_sh = int(height * h_mag)
    width_sh = int(width * w_mag)

    top, left = (top + height) - height_sh, (left + width) - width_sh
    height, width = height_sh, width_sh

    size_list = [top, left, height, width]
    size_list = [pptx.util.Length(i) for i in size_list]
    return size_list


# プレースホルダ分割
def split_size(size_list, h_split, w_split, margin=0):
    result = []
    top, left, height, width = size_list
    mar_i = pptx.util.Cm(margin)
    for i in range(1, h_split + 1, 1):
        for j in range(1, w_split + 1, 1):
            top_s = top + height * ((i - 1) / h_split)
            left_s = left + width * ((j - 1) / w_split)
            height_s = height * (1 / h_split)
            width_s = width * (1 / w_split)
            top_s, left_s, height_s, width_s = top_s + mar_i, left_s + mar_i, height_s - 2 * mar_i, width_s - 2 * mar_i
            split_temp = top_s, left_s, height_s, width_s
            split_temp = [pptx.util.Length(x) for x in split_temp]
            result.append(split_temp)
    return result


# 画像貼り付け
def image_plot(PIL_img, prs, size_list):
    item = save_memory_PIL(PIL_img)
    slide = prs.slides[-1]
    top, left, height, width = size_list
    img = Image.open(item)
    place_aspect_bool, place_aspect = aspect_yokonaga(width, height)
    image_aspect_bool, image_aspect = aspect_yokonaga(*img.size)
    if place_aspect >= image_aspect:
        pic = slide.shapes.add_picture(item, left=left, top=top, width=width)
    else:
        pic = slide.shapes.add_picture(item, left=left, top=top, height=height)
    pic.left = pic.left + pptx.util.Length((width - pic.width) / 2)
    pic.top = pic.top + pptx.util.Length((height - pic.height) / 2)
    return prs


# add table
def df_to_table(prs, df, size_list, font_size, RGB=[0, 0, 0], index=False):
    slide = prs.slides[-1]
    top, left, height, width = size_list
    if index:
        df = df.reset_index()
    row, col = df.shape
    shape = slide.shapes.add_table(rows=row + 1, cols=col, left=left, top=top, width=width, height=height)
    table = shape.table

    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    # columns
    for i, j in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = j

    # value
    for i in range(0, len(df), 1):
        for j in range(0, len(df.columns), 1):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iat[i, j])

    # font_size
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = pptx.util.Pt(font_size)
                run.font.color.rgb = pptx.dml.color.RGBColor(*RGB)

    return prs
