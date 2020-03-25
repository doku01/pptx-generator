from pptx import Presentation
from PIL import Image

def replace_text(replacements, shapes):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        cur_text = paragraph.text
                        new_text = cur_text.replace(str(match), str(replacement))
                        paragraph.text = new_text

if __name__ == '__main__':

    prs = Presentation('assets\\input\\sample.template.pptx')
    # Replace Text
    slides = [slide for slide in prs.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)

    replaces = {
        '{{var1}}': 'text 1',
        '{{var2}}': 'text 2',
        '{{var3}}': 'text 3'
    }
    replace_text(replaces, shapes)

    # Replace imgs
    old_picture = img_shape = prs.slides[1].shapes[4]
    x, y, cx, cy = old_picture.left, old_picture.top, old_picture.width, old_picture.height
    new_picture = prs.slides[1].shapes.add_picture('assets/imgs/inu.jpg', x, y, cx, cy)
    # old_picture.delete()
    # image = Image.open('assets/imgs/inu.jpg', 'r')
    prs.save('assets\\output\\pptx-sample.template.output.pptx')

